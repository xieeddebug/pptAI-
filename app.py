from flask import Flask, request, send_file, render_template, jsonify, make_response, Response
from pptx import Presentation
import os
import tempfile
import traceback
import requests
import json
from docx import Document
from docx.shared import Pt, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.oxml.ns import qn
import io
import base64
import asyncio
import aiohttp
from concurrent.futures import ThreadPoolExecutor
from functools import partial
import re

app = Flask(__name__)

# Dify API配置
DIFY_API_URL = "http://10.119.14.166/v1/chat-messages"
DIFY_API_KEY = "Bearer app-ujLJoBR6bFWdo33nqmgOoEdM"

# PPT对话助手API配置
CHAT_API_URL = "http://10.119.14.166/v1/chat-messages"
CHAT_API_KEY = "Bearer app-aaAtWp12EleFXiKB2L7lgm2J"

def get_dify_response(slide_text):
    headers = {
        "Content-Type": "application/json",
        "Accept": "application/json",
        "Authorization": DIFY_API_KEY
    }
    
    data = {
        "inputs": {},
        "query": f"""请对以下PPT内容进行简单润色和格式调整，要求：
1. 保持原意，不要扩展内容
2. 调整语言更加书面化、严谨
3. 修正明显的语法错误
4. 保持简洁，不要过度发挥
5. 不要生成任何汇报人、日期、时间等信息
6. 如果原文中包含汇报人、日期、时间等信息，请删除这些内容

PPT内容：{slide_text}""",
        "response_mode": "blocking",
        "conversation_id": "",
        "user": "ppt_user"
    }
    
    try:
        print(f"正在调用Dify API，内容：{slide_text[:100]}...")  # 打印前100个字符
        response = requests.post(DIFY_API_URL, headers=headers, json=data)
        print(f"API响应状态码：{response.status_code}")
        print(f"API响应内容：{response.text[:200]}...")  # 打印前200个字符
        
        response.raise_for_status()
        result = response.json()
        
        if 'answer' in result:
            return result['answer']
        elif 'message' in result and 'content' in result['message']:
            return result['message']['content']
        else:
            print(f"API返回结果格式：{json.dumps(result, ensure_ascii=False)}")
            return "无法解析API返回结果"
            
    except requests.exceptions.RequestException as e:
        print(f"API请求错误: {str(e)}")
        return f"API请求错误: {str(e)}"
    except json.JSONDecodeError as e:
        print(f"JSON解析错误: {str(e)}")
        print(f"原始响应内容: {response.text}")
        return "API返回格式错误"
    except Exception as e:
        print(f"其他错误: {str(e)}")
        print(f"错误堆栈: {traceback.format_exc()}")
        return f"发生错误: {str(e)}"

def get_chat_response(prompt):
    headers = {
        "Content-Type": "application/json",
        "Accept": "application/json",
        "Authorization": CHAT_API_KEY
    }
    
    data = {
        "inputs": {},
        "query": prompt,
        "response_mode": "blocking",
        "conversation_id": "",
        "user": "ppt_user"
    }
    
    try:
        print(f"正在调用Chat API，内容：{prompt[:100]}...")  # 打印前100个字符
        response = requests.post(CHAT_API_URL, headers=headers, json=data)
        print(f"API响应状态码：{response.status_code}")
        print(f"API响应内容：{response.text[:200]}...")  # 打印前200个字符
        
        response.raise_for_status()
        result = response.json()
        
        if 'answer' in result:
            return result['answer']
        elif 'message' in result and 'content' in result['message']:
            return result['message']['content']
        else:
            print(f"API返回结果格式：{json.dumps(result, ensure_ascii=False)}")
            return "无法解析API返回结果"
            
    except requests.exceptions.RequestException as e:
        print(f"API请求错误: {str(e)}")
        return f"API请求错误: {str(e)}"
    except json.JSONDecodeError as e:
        print(f"JSON解析错误: {str(e)}")
        print(f"原始响应内容: {response.text}")
        return "API返回格式错误"
    except Exception as e:
        print(f"其他错误: {str(e)}")
        print(f"错误堆栈: {traceback.format_exc()}")
        return f"发生错误: {str(e)}"

@app.route('/')
def index():
    return render_template('index.html')

async def get_dify_response_async(session, slide_text):
    headers = {
        "Content-Type": "application/json",
        "Accept": "application/json",
        "Authorization": DIFY_API_KEY
    }
    
    # 限制文本长度
    max_length = 1500
    if len(slide_text) > max_length:
        # 按句号分割并保留前1500个字符的完整句子
        sentences = slide_text[:max_length].split('。')
        slide_text = '。'.join(sentences[:-1]) + '。'
    
    data = {
        "inputs": {},
        "query": f"""请对以下PPT内容进行简单润色和格式调整，要求：
1. 保持原意，不要扩展内容
5. 不要生成任何汇报人、日期、时间等信息
6. 如果原文中包含汇报人、日期、时间等信息，请删除这些内容

PPT内容：{slide_text}""",
        "response_mode": "blocking",
        "conversation_id": "",
        "user": "ppt_user"
    }
    
    try:
        async with session.post(DIFY_API_URL, headers=headers, json=data) as response:
            if response.status != 200:
                error_text = await response.text()
                print(f"API错误响应: {error_text}")
                return f"API请求失败: {response.status}"
            
            result = await response.json()
            if 'answer' in result:
                return result['answer']
            elif 'message' in result and 'content' in result['message']:
                return result['message']['content']
            else:
                print(f"API返回格式异常: {result}")
                return "无法生成备注"
    except Exception as e:
        print(f"生成备注时出错: {str(e)}")
        return f"生成备注时出错: {str(e)}"

def process_notes_collection(prs):
    doc = Document()
    doc.add_heading('PPT备注合集', 0).alignment = WD_ALIGN_PARAGRAPH.CENTER
    
    for i, slide in enumerate(prs.slides, 1):
        # 添加页码标题
        doc.add_heading(f'第{i}页', level=1).alignment = WD_ALIGN_PARAGRAPH.LEFT
        
        notes_slide = slide.notes_slide
        if notes_slide and notes_slide.notes_text_frame.text:
            notes_text = notes_slide.notes_text_frame.text
            
            # 如果备注内容以"第X页"开头，移除这部分
            notes_text = re.sub(r'^第\d+页[。，：:\n\s]*', '', notes_text).strip()
        else:
            notes_text = "（此页无备注）"
        
        p = doc.add_paragraph()
        p.paragraph_format.line_spacing = 1.5
        p.paragraph_format.first_line_indent = Pt(0)
        run = p.add_run(notes_text)
        run.font.name = '仿宋'
        run._element.rPr.rFonts.set(qn('w:eastAsia'), '仿宋')
        run.font.size = Pt(12)
        
        doc.add_paragraph('', style='Normal')
    
    return doc

def should_skip_text(text):
    # 检查是否包含需要跳过的信息
    if not text:
        return True
        
    # 转换为小写进行检查
    text_lower = text.lower()
    
    # 检查关键词
    skip_keywords = ['汇报人', '单位', 'monday', 'tuesday', 'wednesday', 'thursday', 'friday', 'saturday', 'sunday']
    if any(keyword in text_lower for keyword in skip_keywords):
        return True
    
    # 检查年份（2000-2099）
    if any(str(year) in text for year in range(2000, 2100)):
        return True
    
    # 检查月份
    months = ['月', 'january', 'february', 'march', 'april', 'may', 'june', 'july', 'august', 'september', 'october', 'november', 'december']
    if any(month in text_lower for month in months):
        return True
    
    return False

def extract_text_from_shape(shape):
    text = []  # 使用列表来保持文本顺序
    
    # 处理普通文本框
    if hasattr(shape, "text"):
        shape_text = shape.text.strip()
        if shape_text:
            text.append(shape_text)
    
    # 处理表格
    if hasattr(shape, "has_table") and shape.has_table:
        for row in shape.table.rows:
            row_text = ""
            for cell in row.cells:
                cell_text = cell.text.strip()
                if cell_text:
                    row_text += cell_text + " "
            if row_text:
                text.append(row_text.strip())
    
    # 处理组合形状
    if hasattr(shape, "shapes"):
        for sub_shape in shape.shapes:
            sub_text = extract_text_from_shape(sub_shape)
            text.extend(sub_text)  # 合并子形状的文本列表
    
    # 处理SmartArt
    if hasattr(shape, "graphic_frame"):
        # 处理SmartArt图形
        if hasattr(shape.graphic_frame, "graphic_data"):
            smart_art_text = ""
            for element in shape.graphic_frame.graphic_data.iter():
                if element.tag.endswith('}t'):  # 查找文本元素
                    element_text = element.text.strip() if element.text else ""
                    if element_text:
                        smart_art_text += element_text + " "
            if smart_art_text:
                text.append(smart_art_text.strip())
        
        # 处理图表
        if hasattr(shape.graphic_frame, "chart"):
            chart = shape.graphic_frame.chart
            chart_text = ""
            
            # 提取图表标题
            if hasattr(chart, "has_title") and chart.has_title:
                if hasattr(chart.title, "text_frame") and chart.title.text_frame.text:
                    title = chart.title.text_frame.text.strip()
                    if title:
                        chart_text += title + " "
            
            # 提取图表数据标签
            if hasattr(chart, "plots"):
                for plot in chart.plots:
                    if hasattr(plot, "data_labels"):
                        for label in plot.data_labels:
                            if hasattr(label, "text_frame") and label.text_frame and label.text_frame.text:
                                chart_text += label.text_frame.text.strip() + " "
            
            if chart_text:
                text.append(chart_text.strip())
    
    # 处理文本框架
    if hasattr(shape, "text_frame"):
        try:
            text_frame_content = ""
            # 处理普通文本框架
            if shape.text_frame.text:
                text_frame_content += shape.text_frame.text.strip() + " "
            # 处理段落
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    if run.text:
                        text_frame_content += run.text.strip() + " "
            if text_frame_content:
                text.append(text_frame_content.strip())
        except Exception as e:
            print(f"处理文本框架时出错: {str(e)}")
    
    # 处理占位符
    try:
        if hasattr(shape, "is_placeholder") and shape.is_placeholder:
            if hasattr(shape, "text"):
                placeholder_text = shape.text.strip()
                if placeholder_text:
                    text.append(placeholder_text)
    except Exception as e:
        print(f"处理占位符时出错: {str(e)}")
    
    return text

def extract_slide_text(slide):
    all_text = []  # 使用列表存储所有文本，保持顺序
    
    # 提取所有形状中的文本
    for shape in slide.shapes:
        shape_text = extract_text_from_shape(shape)
        all_text.extend(shape_text)
    
    # 处理页眉页脚
    if hasattr(slide, "header"):
        header_text = slide.header.text.strip()
        if header_text:
            all_text.append(header_text)
    if hasattr(slide, "footer"):
        footer_text = slide.footer.text.strip()
        if footer_text:
            all_text.append(footer_text)
    
    # 处理备注
    if hasattr(slide, "notes_slide") and slide.notes_slide:
        notes_text = slide.notes_slide.notes_text_frame.text.strip()
        if notes_text:
            all_text.append(notes_text)
    
    # 去重并保持顺序
    seen = set()
    unique_text = []
    for text in all_text:
        if text not in seen:
            seen.add(text)
            unique_text.append(text)
    
    # 返回去重后的文本，保持原有顺序
    return "\n".join(unique_text)

@app.route('/api/process-ppt', methods=['POST'])
def process_ppt():
    if 'file' not in request.files:
        return jsonify({'error': '没有上传文件'}), 400
    
    file = request.files['file']
    if file.filename == '':
        return jsonify({'error': '没有选择文件'}), 400
    
    if not file.filename.endswith(('.ppt', '.pptx')):
        return jsonify({'error': '请上传PPT文件'}), 400

    try:
        # 获取任务选项
        generate_ppt_notes = request.form.get('generate_ppt_notes') == 'true'
        generate_notes_collection = request.form.get('generate_notes_collection') == 'true'

        if not generate_ppt_notes and not generate_notes_collection:
            return jsonify({'error': '请至少选择一个生成任务'}), 400

        # 创建临时文件保存上传的PPT
        with tempfile.NamedTemporaryFile(delete=False, suffix='.pptx') as temp_input:
            file.save(temp_input.name)
            input_path = temp_input.name

        # 创建临时文件保存处理后的PPT
        with tempfile.NamedTemporaryFile(delete=False, suffix='.pptx') as temp_output:
            output_path = temp_output.name
        
        # 打开PPT文件
        prs = Presentation(input_path)
        
        # 收集所有幻灯片文本
        slide_texts = []
        for i, slide in enumerate(prs.slides, 1):
            slide_text = extract_slide_text(slide)
            slide_texts.append((i, slide, slide_text))

        # 如果需要生成PPT备注或备注合集，都需要先生成备注
        if generate_ppt_notes or generate_notes_collection:
            async def process_slides():
                async with aiohttp.ClientSession() as session:
                    # 并行处理每页的备注
                    tasks = []
                    for i, slide, slide_text in slide_texts:
                        task = get_dify_response_async(session, slide_text)
                        tasks.append(task)
                    
                    # 等待所有备注生成完成
                    return await asyncio.gather(*tasks)

            # 使用asyncio.run()运行异步任务
            loop = asyncio.new_event_loop()
            asyncio.set_event_loop(loop)
            notes_contents = loop.run_until_complete(process_slides())
            loop.close()

            # 添加备注到PPT
            for (i, slide, _), notes_content in zip(slide_texts, notes_contents):
                if not slide.has_notes_slide:
                    slide.notes_slide
                notes_slide = slide.notes_slide
                notes_text_frame = notes_slide.notes_text_frame
                notes_text_frame.clear()  # 清除现有备注
                notes_text_frame.text = notes_content

        # 保存处理后的PPT
        prs.save(output_path)

        # 获取原始文件名（不含扩展名）
        original_name = os.path.splitext(file.filename)[0]

        # 创建ZIP文件
        with tempfile.NamedTemporaryFile(delete=False, suffix='.zip') as temp_zip:
            import zipfile
            with zipfile.ZipFile(temp_zip.name, 'w') as zipf:
                # 如果生成了PPT备注且需要下载PPT，添加到ZIP
                if generate_ppt_notes:
                    zipf.write(output_path, f'{original_name}（带备注）.pptx')
                
                # 如果需要生成备注合集
                if generate_notes_collection:
                    # 直接生成备注合集文档
                    notes_doc = process_notes_collection(prs)
                    # 保存备注合集文档
                    temp_notes = tempfile.NamedTemporaryFile(delete=False, suffix='.docx')
                    notes_doc.save(temp_notes.name)
                    notes_path = temp_notes.name
                    # 先关闭文件
                    temp_notes.close()
                    # 添加到ZIP
                    zipf.write(notes_path, f'{original_name}（备注合集）.docx')
                    # 删除临时文件
                    try:
                        os.unlink(notes_path)
                    except Exception as e:
                        print(f"删除临时文件失败: {str(e)}")
        
        # 发送ZIP文件
        return send_file(
            temp_zip.name,
            as_attachment=True,
            download_name=f'{original_name}（处理结果）.zip',
            mimetype='application/zip'
        )

    except Exception as e:
        print(f"处理文件时出错: {str(e)}")
        print(f"错误堆栈: {traceback.format_exc()}")
        # 清理临时文件
        if 'input_path' in locals():
            os.unlink(input_path)
        if 'output_path' in locals():
            os.unlink(output_path)
        if 'temp_zip' in locals():
            os.unlink(temp_zip.name)
        return jsonify({'error': str(e)}), 500

@app.route('/api/get-ppt-content', methods=['POST'])
def get_ppt_content():
    if 'file' not in request.files:
        return jsonify({'error': '没有上传文件'}), 400
    
    file = request.files['file']
    if file.filename == '':
        return jsonify({'error': '没有选择文件'}), 400
    
    if not file.filename.endswith(('.ppt', '.pptx')):
        return jsonify({'error': '请上传PPT文件'}), 400

    try:
        # 保存上传的文件
        temp_file = tempfile.NamedTemporaryFile(delete=False, suffix='.pptx')
        file.save(temp_file.name)
        
        # 读取PPT内容
        prs = Presentation(temp_file.name)
        content = []
        
        for slide in prs.slides:
            slide_text = extract_slide_text(slide)
            content.append(slide_text.strip())
        
        # 清理临时文件
        os.unlink(temp_file.name)
        
        return jsonify({'content': '\n\n'.join(content)})
    except Exception as e:
        return jsonify({'error': str(e)}), 500

@app.route('/api/chat', methods=['POST'])
def chat():
    if 'file' not in request.files:
        return jsonify({'error': '没有上传文件'}), 400
    
    file = request.files['file']
    if file.filename == '':
        return jsonify({'error': '没有选择文件'}), 400
    
    if not (file.filename.endswith('.ppt') or file.filename.endswith('.pptx')):
        return jsonify({'error': '请上传PPT文件（.ppt或.pptx格式）'}), 400

    try:
        # 创建临时文件保存上传的PPT
        with tempfile.NamedTemporaryFile(delete=False, suffix='.pptx') as temp_input:
            file.save(temp_input.name)
            input_path = temp_input.name

        # 打开PPT文件
        prs = Presentation(input_path)
        
        # 提取前15页的文本
        all_text = ""
        slides = list(prs.slides)  # 转换为列表
        max_pages = min(15, len(slides))  # 取前15页或总页数的较小值
        
        for i in range(max_pages):
            slide_text = extract_slide_text(slides[i])
            all_text += f"第{i+1}页：\n{slide_text}\n"
            
        # 如果PPT超过15页，添加提示信息
        if len(slides) > 15:
            all_text += "\n注：由于PPT内容较多，仅展示前15页内容作为参考。"

        # 获取用户问题
        data = request.form
        question = data.get('question', '')
        
        # 构建prompt
        prompt = f"""基于以下PPT内容回答问题。如果问题与PPT内容无关，请说明无法回答。

PPT内容：
{all_text}

问题：{question}

请提供准确、简洁的回答。"""
        
        # 调用Chat API
        response = get_chat_response(prompt)
        return jsonify({'answer': response})

    except Exception as e:
        print(f"处理对话时出错: {str(e)}")
        print(f"错误堆栈: {traceback.format_exc()}")
        # 清理临时文件
        if 'input_path' in locals():
            os.unlink(input_path)
        return jsonify({'error': f'处理对话时出错: {str(e)}'}), 500

@app.route('/api/generate-notes-collection', methods=['POST'])
def generate_notes_collection():
    if 'file' not in request.files:
        return jsonify({'error': '没有上传文件'}), 400
    
    file = request.files['file']
    if file.filename == '':
        return jsonify({'error': '没有选择文件'}), 400
    
    if not file.filename.endswith(('.ppt', '.pptx')):
        return jsonify({'error': '请上传PPT文件'}), 400

    try:
        # 创建临时文件保存上传的PPT
        with tempfile.NamedTemporaryFile(delete=False, suffix='.pptx') as temp_input:
            file.save(temp_input.name)
            input_path = temp_input.name

        # 打开PPT文件并添加备注
        prs = Presentation(input_path)
        
        # 为每一页添加备注
        for i, slide in enumerate(prs.slides, 1):
            print(f"\n处理第 {i} 页...")
            
            # 获取幻灯片文本内容
            slide_text = extract_slide_text(slide)
            
            print(f"提取的文本内容：{slide_text[:100]}...")
            
            # 获取Dify生成的备注内容
            notes_content = get_dify_response(slide_text)
            print(f"生成的备注内容：{notes_content}")
            
            if not notes_content:
                notes_content = f"第{i}页备注生成失败"
            
            # 获取或创建备注页
            notes_slide = slide.notes_slide
            notes_text_frame = notes_slide.notes_text_frame
            
            # 添加备注
            notes_text_frame.text = notes_content

        # 创建临时文件保存处理后的PPT
        with tempfile.NamedTemporaryFile(delete=False, suffix='.pptx') as temp_ppt:
            prs.save(temp_ppt.name)
            processed_ppt_path = temp_ppt.name
        
        # 创建Word文档
        doc = Document()
        
        # 设置标题
        title = doc.add_heading('PPT备注合集', 0)
        title.alignment = WD_ALIGN_PARAGRAPH.CENTER
        
        # 收集每页的备注
        for i, slide in enumerate(prs.slides, 1):
            # 添加页码标题
            page_title = doc.add_heading(f'第{i}页备注', level=1)
            page_title.alignment = WD_ALIGN_PARAGRAPH.LEFT
            
            # 获取备注内容
            notes_slide = slide.notes_slide
            if notes_slide and notes_slide.notes_text_frame.text:
                notes_text = notes_slide.notes_text_frame.text
            else:
                notes_text = "（此页无备注）"
            
            # 添加备注内容
            p = doc.add_paragraph()
            p.paragraph_format.line_spacing = 1.5
            p.paragraph_format.first_line_indent = Pt(24)
            run = p.add_run(notes_text)
            run.font.name = '仿宋'
            run._element.rPr.rFonts.set(qn('w:eastAsia'), '仿宋')
            run.font.size = Pt(12)
            
            # 添加分隔行
            doc.add_paragraph('', style='Normal')
        
        # 保存Word文档
        with tempfile.NamedTemporaryFile(delete=False, suffix='.docx') as temp_doc:
            doc.save(temp_doc.name)
            doc_path = temp_doc.name

        # 创建ZIP文件包含两个文件
        with tempfile.NamedTemporaryFile(delete=False, suffix='.zip') as temp_zip:
            import zipfile
            with zipfile.ZipFile(temp_zip.name, 'w') as zipf:
                zipf.write(processed_ppt_path, 'PPT带备注.pptx')
                zipf.write(doc_path, 'PPT备注合集.docx')

        # 发送ZIP文件
        return send_file(
            temp_zip.name,
            as_attachment=True,
            download_name='PPT备注文件.zip',
            mimetype='application/zip'
        )

    except Exception as e:
        print(f"生成备注合集时出错: {str(e)}")
        print(f"错误堆栈: {traceback.format_exc()}")
        # 清理临时文件
        if 'input_path' in locals():
            os.unlink(input_path)
        if 'processed_ppt_path' in locals():
            os.unlink(processed_ppt_path)
        if 'doc_path' in locals():
            os.unlink(doc_path)
        if 'temp_zip' in locals():
            os.unlink(temp_zip.name)
        return jsonify({'error': f'生成备注合集时出错: {str(e)}'}), 500

if __name__ == '__main__':
    import sys
    import platform
    
    if platform.system() == 'Windows':
        # Windows系统使用事件循环策略
        if sys.version_info[0] == 3 and sys.version_info[1] >= 8:
            asyncio.set_event_loop_policy(asyncio.WindowsSelectorEventLoopPolicy())
    
    app.run(debug=True, host='127.0.0.1', port=5000) 